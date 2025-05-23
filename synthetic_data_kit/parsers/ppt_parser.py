# Copyright (c) Meta Platforms, Inc. and affiliates.
# All rights reserved.
#
# This source code is licensed under the terms described in the LICENSE file in
# the root directory of this source tree.
# PPTX parser logic

import os
import lance
import pyarrow as pa
from typing import Dict, Any

from pptx.enum.shapes import MSO_SHAPE_TYPE

class PPTParser:
    """Parser for PowerPoint presentations"""

    def parse(self, file_path: str, multimodal: bool = False) -> any:
        """Parse a PPTX file.
        
        Args:
            file_path: Path to the PPTX file
            multimodal: If True, extract text and images slide by slide.
                        Otherwise, extract all text into a single string.
            
        Returns:
            If multimodal is False, returns a string with all extracted text.
            If multimodal is True, returns a list of dictionaries, where each
            dictionary has 'text' (from a slide) and 'image' (first image
            bytes from that slide, or None).
        """
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("python-pptx is required for PPTX parsing. Install it with: pip install python-pptx")

        prs = Presentation(file_path)
        
        if not multimodal:
            all_text = []
            for i, slide in enumerate(prs.slides):
                slide_text_parts = []
                slide_text_parts.append(f"--- Slide {i+1} ---")
                
                if slide.shapes.title and slide.shapes.title.has_text_frame and slide.shapes.title.text_frame.text:
                    slide_text_parts.append(f"Title: {slide.shapes.title.text_frame.text}")
                
                for shape in slide.shapes:
                    if shape.has_text_frame and shape.text_frame.text:
                        slide_text_parts.append(shape.text_frame.text)
                
                all_text.append("\n".join(slide_text_parts))
            return "\n\n".join(all_text)
        else:
            output_data = []
            for i, slide in enumerate(prs.slides):
                slide_text_parts = []
                slide_text_parts.append(f"--- Slide {i+1} ---")
                
                # Get slide title
                if slide.shapes.title and slide.shapes.title.has_text_frame and slide.shapes.title.text_frame.text:
                    slide_text_parts.append(f"Title: {slide.shapes.title.text_frame.text}")

                # Get text from shapes
                for shape in slide.shapes:
                    if shape.has_text_frame and shape.text_frame.text:
                        slide_text_parts.append(shape.text_frame.text)
                
                slide_text_content = "\n".join(slide_text_parts)

                # Extract first image from the slide
                first_image_bytes = None
                for shape in slide.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            if hasattr(shape, 'image') and hasattr(shape.image, 'blob'):
                                first_image_bytes = shape.image.blob
                                break # Found the first image
                        except Exception:
                             # In case of issues accessing blob, ensure it's None
                            first_image_bytes = None
                            break 
                
                output_data.append({
                    'text': slide_text_content,
                    'image': first_image_bytes
                })
            return output_data

    def save(self, content: any, output_path: str) -> None:
        """Save the extracted content to a Lance file.
        
        Args:
            content: Extracted content (string or list of dicts)
            output_path: Path to save the Lance file
        """
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        
        if isinstance(content, str):
            # Text-only mode
            data = [pa.array([content])]
            names = ['text']
            table = pa.Table.from_arrays(data, names=names)
        elif isinstance(content, list):
            # Multimodal mode
            texts = pa.array([item['text'] for item in content], type=pa.string())
            images_data = [item.get('image') for item in content] # Handles None if 'image' key is missing too
            images = pa.array(images_data, type=pa.binary())
            
            schema = pa.schema([
                ('text', pa.string()),
                ('image', pa.binary())
            ])
            table = pa.Table.from_arrays([texts, images], schema=schema)
        else:
            raise ValueError("Unsupported content type for saving.")
            
        lance.write_dataset(table, output_path, mode="overwrite")